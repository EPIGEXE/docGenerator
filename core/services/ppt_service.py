from pptx import Presentation
from pptx.util import Inches, Pt
import os

class PPTService:
    async def generate_ppt(self, analyzed_content: dict) -> str:
        """PPT 파일 생성"""
        try:
            prs = Presentation()
            
            # 제목 슬라이드 추가
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            title.text = analyzed_content.get('title', '프로젝트 명세서')
            
            # 섹션별 슬라이드 생성
            for section in analyzed_content.get('sections', []):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                # 섹션 제목
                title = slide.shapes.title
                title.text = section.get('title', '')
                
                # 기능 목록
                content = slide.placeholders[1]
                tf = content.text_frame
                
                for feature in section.get('features', []):
                    p = tf.add_paragraph()
                    p.text = f"• {feature.get('name', '')}"
                    
                    # 하위 기능 목록
                    for sub in feature.get('sub_features', []):
                        p = tf.add_paragraph()
                        p.text = f"    - {sub}"
                        p.level = 1
            
            # 파일 저장
            output_path = os.path.join('media', 'documents', 'temp.pptx')
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            prs.save(output_path)
            
            return output_path
            
        except Exception as e:
            raise Exception(f"PPT 생성 중 오류 발생: {str(e)}")